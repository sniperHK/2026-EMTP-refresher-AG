#!/usr/bin/env python3
"""
generate_pptx_v2.py — 產生 M01 & M02 課程簡報 PPTX
2026 台北市 EMTP 複訓
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x2D, 0x5F, 0x8A)
LIGHT_BLUE_HDR = RGBColor(0xD6, 0xE8, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_TEXT = RGBColor(0x1A, 0x1A, 0x1A)
ROW_ALT_1 = RGBColor(0xF5, 0xF9, 0xFC)
ROW_ALT_2 = RGBColor(0xE8, 0xF0, 0xF7)
ACCENT_RED = RGBColor(0xCC, 0x33, 0x33)

FONT_NAME = "Calibri"

# ── Slide dimensions (16:9) ─────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Output path ─────────────────────────────────────────────
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
PROJECT_ROOT = os.path.dirname(SCRIPT_DIR)
OUTPUT_DIR = os.path.join(PROJECT_ROOT, "exports", "2026-02-26_v2")


# ════════════════════════════════════════════════════════════
#  Helper functions
# ════════════════════════════════════════════════════════════

def new_presentation():
    """Return a blank 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def _set_font(run, size=Pt(18), bold=False, color=DARK_TEXT, name=FONT_NAME):
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name


def _fill_bg(slide, color):
    """Set slide background to a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_slide(prs, title_text, subtitle_text):
    """Dark blue background title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _fill_bg(slide, DARK_BLUE)

    # Title
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(2.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title_text
    _set_font(run, size=Pt(36), bold=True, color=WHITE)

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(4.2), Inches(11.7), Inches(1.2))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = subtitle_text
    _set_font(run2, size=Pt(22), bold=False, color=WHITE)

    return slide


def add_section_header(prs, text):
    """Medium blue background section header slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_bg(slide, MEDIUM_BLUE)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(11.7), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    _set_font(run, size=Pt(28), bold=True, color=WHITE)

    return slide


def _add_slide_title(slide, text, size=Pt(28)):
    """Add a top-area title textbox."""
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.0), Inches(0.9))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    _set_font(run, size=size, bold=True, color=DARK_BLUE)
    return tf


def add_content_slide(prs, title, bullets, *, numbered=False, note_text=None):
    """White slide with a title and bullet/numbered list."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, title)

    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = f"{i + 1}. " if numbered else "\u2022 "
        run = p.add_run()
        run.text = prefix + bullet
        _set_font(run, size=Pt(18), color=DARK_TEXT)
        p.space_after = Pt(6)

    if note_text:
        p = tf.add_paragraph()
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = note_text
        _set_font(run, size=Pt(14), bold=True, color=ACCENT_RED)

    return slide


def add_quote_slide(prs, title, quote_text):
    """White slide with a single prominent quote/paragraph."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, title)

    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.0), Inches(4.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = quote_text
    _set_font(run, size=Pt(20), color=DARK_TEXT)
    p.line_spacing = Pt(30)

    return slide


def add_table_slide(prs, title, headers, rows, *, note_text=None, col_widths=None):
    """White slide with a styled table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_slide_title(slide, title)

    n_rows = len(rows) + 1
    n_cols = len(headers)

    # Default equal column widths filling most of the slide
    table_width = Inches(11.5)
    if col_widths is None:
        col_widths = [table_width // n_cols] * n_cols

    table_top = Inches(1.5)
    table_left = Inches(0.9)

    shape = slide.shapes.add_table(n_rows, n_cols, table_left, table_top,
                                   table_width, Inches(0.45 * n_rows))
    table = shape.table

    # Set column widths
    for ci, w in enumerate(col_widths):
        table.columns[ci].width = w

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        _style_cell(cell, bold=True, bg=LIGHT_BLUE_HDR, font_color=DARK_BLUE, size=Pt(16))

    # Data rows
    for ri, row in enumerate(rows):
        bg = ROW_ALT_1 if ri % 2 == 0 else ROW_ALT_2
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            _style_cell(cell, bg=bg, size=Pt(15))

    if note_text:
        # Calculate approximate bottom of table
        note_top = table_top + Inches(0.45 * n_rows) + Inches(0.2)
        txBox = slide.shapes.add_textbox(Inches(0.9), note_top, Inches(11.5), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = note_text
        _set_font(run, size=Pt(14), bold=True, color=ACCENT_RED)

    return slide


def _style_cell(cell, bold=False, bg=None, font_color=DARK_TEXT, size=Pt(15)):
    """Style a single table cell."""
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            _set_font(run, size=size, bold=bold, color=font_color)
    if bg:
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = bg


# ════════════════════════════════════════════════════════════
#  M01 — ALS 病生理學
# ════════════════════════════════════════════════════════════

def build_m01():
    prs = new_presentation()

    # ── Title ──
    add_title_slide(prs,
                    "M01 ALS 病生理學：決策導向的瀕死生理",
                    "2026 台北市 EMTP 複訓")

    # ── §1 核心觀念 ──
    # Slide 2
    add_quote_slide(prs,
                    "核心觀念：生命徵象是生理代償的結果",
                    "當你看到病人 BP 80/40, HR 130 時，你不只是看到「休克」，"
                    "你是看到身體正在「拚命代償以維持灌流」的結果。")

    # Slide 3
    add_table_slide(prs,
                    "瀕死生理的四大支柱",
                    ["支柱", "監測指標", "處置方向"],
                    [
                        ["氧合 Oxygenation", "SpO2", "給氧 (O2 therapy)"],
                        ["通氣 Ventilation", "EtCO2, RR", "輔助呼吸 (BVM/ETT)"],
                        ["灌流 Perfusion", "BP, HR, CRT, 意識", "輸液/升壓藥"],
                        ["代謝 Metabolism", "pH, Base Excess", "矯正根本原因"],
                    ])

    # Slide 4
    add_content_slide(prs,
                      "四大支柱 — 臨床整合",
                      [
                          "四支柱不是獨立存在 — 互相連動",
                          "休克(灌流↓) → 組織缺氧(氧合↓) → 無氧代謝(代謝↓) → 代償性過度換氣(通氣↑)",
                          "處置策略：找到「最先壞掉的那根柱子」修復",
                      ])

    # ── §2 氧合 vs 通氣 ──
    # Slide 5
    add_section_header(prs, "Section 2: 氧合 vs 通氣")

    # Slide 6
    add_table_slide(prs,
                    "氧合 (Oxygenation) vs 通氣 (Ventilation)",
                    ["比較項目", "氧合 Oxygenation", "通氣 Ventilation"],
                    [
                        ["定義", "O2 進入血液", "CO2 排出"],
                        ["主要指標", "SpO2, PaO2", "EtCO2, PaCO2, RR"],
                        ["常見問題", "Hypoxemia", "Hypercapnia"],
                        ["處置", "給氧 (NRM)", "通氣支持 (BVM/ETT)"],
                    ])

    # Slide 7
    add_content_slide(prs,
                      "臨床決策點",
                      [
                          "SpO2↓ + RR 快 → Oxygenation 問題 → 高濃度氧",
                          "SpO2↓ + RR 慢/淺 → Ventilation 問題 → BVM/插管",
                          "SpO2 正常但意識差 → CO2 蓄積 → 看 EtCO2",
                      ])

    # Slide 8
    add_content_slide(prs,
                      "氧氣解離曲線 (Oxyhemoglobin Dissociation Curve)",
                      [
                          "右移（酸/熱/高 CO2）→ Hb 釋放 O2 ↑ → 有利組織",
                          "左移（鹼/低溫/低 CO2）→ Hb 抓緊 O2 → 組織缺氧",
                          "臨床意義：發燒+酸中毒的病人，SpO2 88% 不等於危急",
                      ])

    # Slide 9
    add_content_slide(prs,
                      "EtCO2 的臨床價值",
                      [
                          "正常值：35-45 mmHg",
                          "EtCO2 驟降 → 心輸出量↓ 或心搏停止",
                          "休克/DKA 病人 EtCO2 偏低 = 代償，不要刻意拉回正常",
                      ])

    # ── §3 循環與灌流 ──
    # Slide 10
    add_section_header(prs, "Section 3: 循環與灌流 — 休克的動態進程")

    # Slide 11
    add_table_slide(prs,
                    "休克三階段",
                    ["階段", "機制", "徵象", "處置"],
                    [
                        ["代償期", "交感興奮", "HR↑, SVR↑, BP 正常", "黃金期！早期介入"],
                        ["失代償期", "代償耗竭", "BP↓, HR↑↑, 意識差", "酸中毒→升壓藥效差"],
                        ["不可逆期", "多器官衰竭", "即使 BP 回來仍死亡", "治療已無效"],
                    ])

    # Slide 12
    add_content_slide(prs,
                      "灌流評估三角 Pump-Pipe-Tank",
                      [
                          "Pump (心臟): Rate, Rhythm, Contractility → 心因性休克",
                          "Pipe (血管): SVR, CRT, 皮膚溫度 → 分布性/阻塞性休克",
                          "Tank (容量): JVD, 黏膜, 肺音 → 低血容性休克",
                      ])

    # Slide 13
    add_table_slide(prs,
                    "Tank 空 vs Tank 滿 vs Pump 不力 — 鑑別",
                    ["鑑別", "Tank 空", "Tank 滿/Pump 不力"],
                    [
                        ["JVD", "(-)", "(+)"],
                        ["肺音", "乾淨", "雙側囉音"],
                        ["處置", "輸液", "不能輸液！升壓/強心"],
                    ])

    # Slide 14
    add_content_slide(prs,
                      "鑑別決策樹：同樣低血壓，該輸液還是升壓？",
                      [
                          "SBP < 90 + 皮膚冰冷 + JVD(-) + 肺音乾淨 → Tank 空 → 輸液 (S02)",
                          "SBP < 90 + 皮膚溫暖 + JVD(-) + ±wheeze → Pipe 鬆 → Epi (S03)",
                          "SBP < 90 + 皮膚冰冷 + JVD(+) + 肺音乾淨 → Tank 被擋 → 減壓 (S01)",
                          "SBP < 90 + 皮膚冰冷 + JVD(±) + 雙側囉音 → Pump 不力 → 升壓/強心 (S05)",
                      ])

    # ── §3.5 心因性休克 + APE ──
    # Slide 15
    add_section_header(prs, "Section 3.5: 心因性休克與急性肺水腫")

    # Slide 16
    add_content_slide(prs,
                      "心因性休克：Pump 壞了",
                      [
                          "Forward failure（前向衰竭）：心輸出量↓ → 全身低灌流 → Cold",
                          "Backward failure（後向衰竭）：血液回堵肺 → 肺水腫 → Wet",
                          "Cold + Wet = 最難處理的急重症",
                          "常見成因：AMI（最常見）、高血壓危象、心律不整",
                      ])

    # Slide 17
    add_content_slide(prs,
                      "急性肺水腫 & SCAPE",
                      [
                          "左心衰竭 → 肺靜脈壓↑ → 超過 Starling 平衡 → 血漿滲入肺泡",
                          "SCAPE = 交感風暴 → 血壓飆高 → 急性肺水腫",
                          "特徵：高血壓 + 呼吸窘迫 + 粉紅泡沫痰",
                          "處置：NTG + CPAP/BiPAP + 坐姿（不要急著插管！）",
                      ])

    # Slide 18  — Nohria-Stevenson 矩陣
    add_table_slide(prs,
                    "Nohria-Stevenson 矩陣",
                    ["", "Warm（溫暖）", "Cold（冰冷）"],
                    [
                        ["Dry（乾）", "A: 穩定 → 觀察", "L: 低灌流 → 輸液試探"],
                        ["Wet（濕）", "B: 肺水腫 → NTG + NIV", "C: 心因性休克 → 升壓+強心"],
                    ])

    # Slide 19
    add_table_slide(prs,
                    "Nohria-Stevenson 處置方向",
                    ["象限", "臨床表現", "處置方向"],
                    [
                        ["A Warm+Dry", "穩定，VS 正常", "觀察"],
                        ["B Warm+Wet", "囉音+SpO2↓, 四肢溫暖, BP↑", "NTG + CPAP (SCAPE)"],
                        ["L Cold+Dry", "低 BP+末梢冰冷, 肺音乾淨", "小量輸液試探"],
                        ["C Cold+Wet", "低 BP+末梢冰冷+囉音", "最棘手：升壓+強心"],
                    ])

    # Slide 20
    add_content_slide(prs,
                      "PseudoPEA — 休克的終點站",
                      [
                          "定義：心臟還在跳，但 SBP < 40-50 → 脈搏摸不到",
                          "不是 cardiac arrest，是 extreme shock",
                          "42-86% 的 PEA 其實是 pseudoPEA",
                          "True PEA 存活 < 1%；PseudoPEA ROSC 70-94%",
                          "POCUS 是唯一可靠鑑別",
                      ])

    # Slide 21
    add_table_slide(prs,
                    "True PEA vs PseudoPEA 處置差異",
                    ["項目", "True PEA", "PseudoPEA"],
                    [
                        ["心臟", "無有效收縮", "有收縮但極弱"],
                        ["處置", "CPR + Epi 1 mg", "Push-dose epi 10-20 mcg"],
                        ["CPR", "持續壓胸", "考慮暫停，改升壓"],
                        ["預後", "< 1%", "ROSC 70%+"],
                    ])

    # ── §4 酸鹼與代謝 ──
    # Slide 22
    add_content_slide(prs,
                      "酸鹼與代謝 (EtCO2)",
                      [
                          "代謝性酸中毒 → Kussmaul breathing → EtCO2 偏低 = 代償",
                          "不要用 BVM 把 EtCO2 拉回 35-45 → 會破壞代償 → pH 驟降",
                          "DKA 病人插管：必須匹配原本的高通氣量",
                      ])

    # ── §5 藥物生理效應 ──
    # Slide 23
    add_section_header(prs, "Section 5: 常見 ALS 藥物的生理效應")

    # Slide 24
    add_table_slide(prs,
                    "藥物與處置的生理效應",
                    ["藥物", "作用目標", "預期效應", "風險"],
                    [
                        ["Fluids", "Tank (Preload)", "SV↑ → CO↑ → BP↑", "Cold+Wet 禁大量輸液"],
                        ["Epinephrine", "α1 + β1", "SVR↑, 收縮力↑", "心肌氧耗↑, 心律不整"],
                        ["NTG", "靜脈/動脈擴張", "減前負荷/後負荷", "SBP<90 禁用"],
                        ["CPAP/BiPAP", "氧合 + 前負荷", "肺泡撐開 + 減回心量", "肺水腫首選！"],
                    ])

    # Slide 25
    add_content_slide(prs,
                      "Red Flags — 關鍵決策示警",
                      [
                          "插管前低血壓 → PPV 會進一步降 BP → 先穩循環再插管",
                          "酸中毒病人插管 → 接管後通氣量不足 → pH 驟降 → Cardiac Arrest",
                          "HOP 原則：Hypotension, Oxygenation, pH — 先處理再插管",
                      ])

    # ── §6 模組總結 ──
    # Slide 26
    add_content_slide(prs,
                      "M01 模組總結（6 大重點）",
                      [
                          "區分 Oxygenation vs Ventilation 問題",
                          "休克識別在血壓掉之前 — 看 HR, CRT, 皮膚, 意識",
                          "Pump-Pipe-Tank 定位休克類型 → 處置方向完全不同",
                          "Nohria-Stevenson 分四象限 — Warm+Wet (NTG+NIV) vs Cold+Wet (升壓+強心)",
                          "PseudoPEA ≠ 心搏停止 — 當重度休克治療",
                          "Resuscitation 優先於 Intubation — 肺水腫先 CPAP，不急著插管",
                      ], numbered=True)

    # Slide 27
    add_table_slide(prs,
                    "五個情境的 Pump-Pipe-Tank 對照",
                    ["情境", "壞掉的是", "怎麼修"],
                    [
                        ["S01 張力性氣胸", "Tank 被擋", "減壓"],
                        ["S02 低血容性休克", "Tank 空", "止血 + 輸液"],
                        ["S03 過敏性休克", "Pipe 鬆", "Epi"],
                        ["S05 心因性休克", "Pump 不力", "升壓/強心"],
                        ["S04 OHCA", "全部崩潰", "CPR + 電擊 + 找原因"],
                    ])

    return prs


# ════════════════════════════════════════════════════════════
#  M02 — 藥物動力學/藥效學 (PK/PD)
# ════════════════════════════════════════════════════════════

def build_m02():
    prs = new_presentation()

    # ── Title ──
    add_title_slide(prs,
                    "M02 藥物動力學/藥效學 (PK/PD)",
                    "掌握「時間」與「情境」| 2026 台北市 EMTP 複訓")

    # ── §1 核心觀念 ──
    # Slide 2
    add_quote_slide(prs,
                    "核心觀念：時間軸就是生命線",
                    "藥物打進去不是馬上有效。\n"
                    "Epinephrine 1mg IV → 幾秒後血壓上升？\n"
                    "Midazolam 5mg IM → 幾分鐘後鎮靜？")

    # Slide 3
    add_content_slide(prs,
                      "PK/PD 四大步驟",
                      [
                          "吸收 Absorption: IV/IO = 100% 生物利用率；IM 在休克時極差",
                          "分布 Distribution: 休克→中央循環優先→小心相對過量",
                          "代謝 Metabolism: 肝功能↓ → Duration 延長",
                          "排泄 Excretion: 腎功能↓ → 藥物蓄積",
                      ])

    # Slide 4
    add_content_slide(prs,
                      "IM 在休克中的致命陷阱",
                      [
                          "休克 → 周邊血管收縮 → 肌肉血流↓ → IM 吸收極差",
                          "藥物積在肌肉裡，等休克改善血流恢復 → 一次釋放 = Late Dump",
                          "結論：休克病人禁用 IM！改 IV/IO",
                      ])

    # ── §2 給藥途徑 ──
    # Slide 5
    add_section_header(prs, "Section 2: 給藥途徑實戰")

    # Slide 6
    add_table_slide(prs,
                    "IV vs IO vs IM vs IN 比較",
                    ["途徑", "Onset", "適用", "休克注意"],
                    [
                        ["IV", "10-30 sec", "ALS 黃金標準", "推藥後 Flush"],
                        ["IO", "30-60 sec", "休克/OHCA 首選", "Humeral > Tibial"],
                        ["IM", "5-15 min", "躁動/無法 IV", "休克禁用！Late Dump"],
                        ["IN", "3-5 min", "兒童/止痛", "每孔 < 1 mL"],
                    ])

    # Slide 7
    add_content_slide(prs,
                      "IO 的臨床重點",
                      [
                          "Humeral IO 到達心臟速度 ≈ Central line",
                          "Tibial IO 需加壓輸液",
                          "清醒病人 IO 很痛 → 先 Lidocaine → Flush 衝開",
                          "建立後每次給藥都要 20 mL NS flush",
                      ])

    # Slide 8
    add_content_slide(prs,
                      "Circulation Time 的重要性",
                      [
                          "CPR 時血流極慢 → 不 flush 藥物到不了心臟",
                          "IV/IO 推藥後：20 mL NS flush + 抬高肢體",
                          "Adenosine 半衰期 < 10 秒 → 慢推就沒效",
                      ])

    # Slide 9
    add_content_slide(prs,
                      "給藥途徑選擇流程",
                      [
                          "穩定病人 → IV 首選",
                          "休克/OHCA + IV 困難 → IO（優先 Humeral）",
                          "躁動 + 無 IV → IM（確認非休克狀態）",
                          "兒童止痛/鎮靜 → IN",
                      ])

    # ── §3 特殊情境 PK/PD ──
    # Slide 10
    add_section_header(prs, "Section 3: 特殊情境 PK/PD 調整")

    # Slide 11
    add_content_slide(prs,
                      "休克 / 低灌流",
                      [
                          "Sedation: Start Low, Go Slow → 減量",
                          "Paralytics: Onset 延遲（circulation time 變長）→ 別急著補藥",
                          "Pressors: 受體已被內源性 catecholamine 活化 → 反應可能劇烈",
                      ])

    # Slide 12
    add_content_slide(prs,
                      "心搏停止",
                      [
                          "完全沒循環 → 全靠 CPR 推藥",
                          "IV/IO 後務必 Flush 20 mL + 抬高",
                          "Epi 累積：ROSC 時體內可能積了 3-5 mg → 血壓飆高/心律不整",
                      ])

    # Slide 13
    add_content_slide(prs,
                      "老年人 / 低體溫",
                      [
                          "老人：脂溶性藥物蓄積、肝腎衰退 → Midazolam/Fentanyl 減量 50%",
                          "低體溫：酵素活性↓ → 藥物代謝超慢 → ROSC 後別急著追加鎮靜",
                          "Start Low, Go Slow 原則",
                      ])

    # ── §4 藥物速查表 ──
    # Slide 14
    add_section_header(prs, "Section 4: 藥物速查表（ACLS 2020）")

    # Slide 15 — Epinephrine
    add_table_slide(prs,
                    "Epinephrine — 同一支藥、兩種用法",
                    ["情境", "濃度/泡製", "劑量", "途徑"],
                    [
                        ["Cardiac Arrest", "1:1,000 原液一支", "1 mg q3-5min", "IV/IO"],
                        ["Anaphylaxis", "1:1,000 原液", "0.3-0.5 mg q5-15min", "IM"],
                        ["Push-dose (低 BP)", "二階段稀釋→10mcg/mL", "10-20 mcg q1-2min", "IV"],
                    ],
                    note_text="台灣只有 1:1,000 (1mg/1mL)。心停用原液一支直接 IV/IO；過敏用原液 IM；push-dose 需稀釋")

    # Slide 16 — Amiodarone
    add_content_slide(prs,
                      "Amiodarone",
                      [
                          "OHCA VF/pVT: 1st 300 mg bolus → 2nd 150 mg",
                          "非 arrest VT: 150 mg IV over 10-15 min（慢滴！）",
                          "陷阱：推太快 → 低血壓（血管擴張+心肌抑制）",
                          "記憶口訣：「活人慢滴，死人才 Push」",
                      ])

    # Slide 17 — Adenosine
    add_content_slide(prs,
                      "Adenosine",
                      [
                          "Stable SVT: 6 mg rapid push → 12 mg × 2",
                          "半衰期 < 10 秒 → 必須 rapid push + flush",
                          "打太慢 = 沒效（被 RBC 代謝光了）",
                          "禁用：irregular wide-complex tachycardia",
                      ])

    # Slide 18 — NTG
    add_table_slide(prs,
                    "Nitroglycerin (NTG) — 肺水腫利器",
                    ["項目", "內容"],
                    [
                        ["適應症", "SCAPE / Warm+Wet 肺水腫"],
                        ["劑量", "SL 0.4 mg q3-5min; SCAPE 可 SL 0.8 mg（急診文獻建議，非標準 AHA 劑量，需醫療指導授權）"],
                        ["IV drip", "10-200 mcg/min"],
                        ["Onset", "SL 1-3 min; IV < 1 min"],
                        ["禁忌", "SBP < 90、RV infarct、近 24-48 hr 內使用 PDE-5 抑制劑（Sildenafil/Tadalafil）"],
                    ],
                    note_text="作用：低劑量→靜脈擴張(減前負荷)；高劑量→動脈擴張(減後負荷)")

    # Slide 19 — NTG PK/PD
    add_content_slide(prs,
                      "NTG 的 PK/PD 教學重點",
                      [
                          "NTG SL 在低灌流時吸收差 — 跟 IM 在休克時失效的道理一樣",
                          "黏膜血流不足 → 吸收延遲 → 要考慮 IV drip",
                          "從 Warm+Wet 惡化到 Cold+Wet → 立即停 NTG！",
                          "SBP < 90 是硬門檻 — 不可妥協",
                      ])

    # Slide 20 — Furosemide & Dobutamine
    add_table_slide(prs,
                    "Furosemide & Dobutamine（速查，不口頭講解）",
                    ["藥物", "適應症", "劑量", "陷阱"],
                    [
                        ["Furosemide", "Volume overload", "20-80 mg IV", "SCAPE ≠ volume overload → 非首選"],
                        ["Dobutamine", "Cold+Wet, SBP>70", "2-20 mcg/kg/min", "增加心肌氧耗；SBP<70 不單獨用"],
                    ],
                    note_text="放學員講義速查表，課堂不佔時間")

    # Slide 21 — Midazolam & Fentanyl
    add_content_slide(prs,
                      "Midazolam & Fentanyl",
                      [
                          "Midazolam — Seizure: IM 10mg / IV 5mg / IN 0.2mg/kg；Sedation: IV 2-5mg / IM 5mg",
                          "老人減量 50%；呼吸抑制！",
                          "Fentanyl: IV 1-2 mcg/kg / IN 1.5 mcg/kg",
                          "不釋放組織胺 → 比 Morphine 更不影響 BP",
                      ])

    # Slide 22 — Norepinephrine
    add_content_slide(prs,
                      "Norepinephrine",
                      [
                          "Septic/Distributive Shock: 0.01-0.5 mcg/kg/min IV drip（從低劑量 titrate）",
                          "Cardiogenic Shock (Cold+Wet, SBP<70): 首選升壓",
                          "外滲 → 組織壞死 → 優先 central line 或大口徑 IV",
                          "泡製：4 mg / 250 mL D5W",
                      ])

    # Slide 23 — Albuterol
    add_content_slide(prs,
                      "Salbutamol (= Albuterol)",
                      [
                          "Bronchospasm: 2.5 mg nebulizer q15-20min",
                          "Inhalation onset 5-15 min",
                          "Side effect: Tachycardia, Tremor",
                      ])

    # ── §5 藥物計算 ──
    # Slide 24
    add_section_header(prs, "Section 5: 藥物計算與泡製")

    # Slide 25
    add_content_slide(prs,
                      "Push-dose Epinephrine 泡製",
                      [
                          "Step 1: 取 Epi 1:1,000 (1mg/1mL) 抽 1mL + NS 至 10mL = 100 mcg/mL",
                          "Step 2: 從 Step 1 取 1mL (100mcg) + NS 至 10mL = 10 mcg/mL",
                          "給藥: 每次 1-2 mL (10-20 mcg) IV push",
                          "Titrate 至 SBP > 90，每 1-2 分鐘可重複",
                      ])

    # Slide 26 — PseudoPEA dosing
    add_content_slide(prs,
                      "PseudoPEA — 為什麼不能用 ACLS 劑量",
                      [
                          "PseudoPEA 的心臟還在跳 → 不是 cardiac arrest",
                          "ACLS Epi 1 mg = push-dose 的 50-100 倍",
                          "全量打進還在跳的心臟 → 心肌過度興奮 → VF/VT",
                          "正確做法：push-dose 10-20 mcg IV q1-2min",
                          "記住：「會跳的心用小量，不跳的心用大量」",
                      ])

    # ── §6 模組總結 ──
    # Slide 27
    add_content_slide(prs,
                      "M02 模組總結（6 大重點）",
                      [
                          "休克病人不要 IM — 改 IV/IO",
                          "Humeral IO 快於 Tibial IO",
                          "老人鎮靜藥減量 50%",
                          "Push-dose Epi 是處理休克的重要技能",
                          "NTG 是肺水腫利器 — SBP<90 或右心梗禁用",
                          "PseudoPEA 用 push-dose (10-20 mcg)，不是 ACLS 劑量 (1 mg)",
                      ], numbered=True)

    # Slide 28
    add_content_slide(prs,
                      "高風險錯誤防呆 — 4 條",
                      [
                          "休克打 IM: 血流不到肌肉 → 改用 IV/IO",
                          "追藥太快: 尊重 circulation time, 至少等 3 min",
                          "IO 忘 Flush: 每次給藥後 20 mL NS flush + 抬高",
                          "Amiodarone 活人快推: 非 arrest → 10-15 min 慢滴",
                      ], numbered=True)

    # Slide 29 — Protocol 聲明
    add_quote_slide(prs,
                    "Protocol 聲明",
                    "藥物劑量依據 AHA ACLS 2020 通用準則\n"
                    "實際執行請以台北市醫療指導醫師核准之在地 Protocol 為準\n\n"
                    "2026 台北市 EMTP 複訓 | ALS 病生理學 & 藥物動力學")

    return prs


# ════════════════════════════════════════════════════════════
#  Main
# ════════════════════════════════════════════════════════════

def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    # Build and save M01
    m01 = build_m01()
    m01_path = os.path.join(OUTPUT_DIR, "M01_ALS病生理學_v2.pptx")
    m01.save(m01_path)
    print(f"[OK] M01 saved → {m01_path}  ({len(m01.slides)} slides)")

    # Build and save M02
    m02 = build_m02()
    m02_path = os.path.join(OUTPUT_DIR, "M02_藥物動力學PK-PD_v2.pptx")
    m02.save(m02_path)
    print(f"[OK] M02 saved → {m02_path}  ({len(m02.slides)} slides)")


if __name__ == "__main__":
    main()
